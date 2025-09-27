from pptx import Presentation

def generate_presentation(content: str, filename="output.pptx"):
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]  # title layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hackathon Demo"
    subtitle.text = content

    # Save file
    prs.save(filename)
    return filename
